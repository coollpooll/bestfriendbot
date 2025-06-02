import os
import logging
import re
import datetime
import asyncio
from io import BytesIO
from aiogram import Bot, Dispatcher, types, F
from aiogram.fsm.storage.memory import MemoryStorage
from aiogram.filters import CommandStart
from aiogram.enums import ParseMode
from aiogram.client.bot import DefaultBotProperties
from dotenv import load_dotenv
import asyncpg
from fastapi import FastAPI, Request
from openai import OpenAI
from pydub import AudioSegment
from aiogram.types import ReplyKeyboardMarkup, KeyboardButton

# Документы
import mimetypes
import zipfile
import rarfile
from PyPDF2 import PdfReader
from docx import Document as DocxDocument
import csv
import pandas as pd
import pptx
import xlrd
import openpyxl
import textract

load_dotenv()
logging.basicConfig(level=logging.INFO)

BOT_TOKEN = os.getenv("BOT_TOKEN")
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")
ASSISTANT_ID = os.getenv("ASSISTANT_ID")
DATABASE_URL = os.getenv("DATABASE_URL")
OWNER_CHAT_ID = int(os.getenv("OWNER_CHAT_ID"))

bot = Bot(token=BOT_TOKEN, default=DefaultBotProperties(parse_mode=ParseMode.HTML))
dp = Dispatcher(storage=MemoryStorage())
app = FastAPI()
openai_client = OpenAI(api_key=OPENAI_API_KEY)

# --- Database logic
class Database:
    def __init__(self, dsn):
        self.dsn = dsn
        self.pool = None

    async def connect(self):
        self.pool = await asyncpg.create_pool(self.dsn)

    async def disconnect(self):
        if self.pool:
            await self.pool.close()

    async def add_user(self, user_id):
        async with self.pool.acquire() as connection:
            await connection.execute(
                "INSERT INTO users (user_id, requests_today) VALUES ($1, $2) ON CONFLICT (user_id) DO NOTHING",
                user_id, 0
            )

    async def add_message(self, user_id, role, content):
        async with self.pool.acquire() as connection:
            await connection.execute(
                "INSERT INTO dialog_history (user_id, role, content) VALUES ($1, $2, $3)",
                user_id, role, content
            )

    async def get_history(self, user_id, limit=16):
        async with self.pool.acquire() as connection:
            rows = await connection.fetch(
                """
                SELECT role, content FROM dialog_history
                WHERE user_id = $1
                ORDER BY created_at DESC
                LIMIT $2
                """,
                user_id, limit
            )
            return [{"role": row["role"], "content": row["content"]} for row in reversed(rows)]

    async def add_subscription(self, user_id, plan, payment_id):
        async with self.pool.acquire() as connection:
            if plan == 'monthly':
                expires_at = "NOW() + INTERVAL '30 days'"
            elif plan == 'yearly':
                expires_at = "NOW() + INTERVAL '365 days'"
            else:
                raise ValueError("Unknown plan type")
            await connection.execute(
                f"""
                INSERT INTO subscriptions 
                    (user_id, plan, started_at, expires_at, payment_id, status) 
                VALUES 
                    ($1, $2, NOW(), {expires_at}, $3, 'active')
                """,
                user_id, plan, payment_id
            )

    async def get_user_subscription(self, user_id):
        async with self.pool.acquire() as connection:
            row = await connection.fetchrow(
                """
                SELECT plan, expires_at FROM subscriptions
                WHERE user_id = $1 AND status = 'active' AND expires_at > NOW()
                ORDER BY expires_at DESC
                LIMIT 1
                """,
                user_id
            )
            return row

    async def get_stats(self):
        async with self.pool.acquire() as connection:
            users = await connection.fetchval("SELECT COUNT(*) FROM users")
            monthly = await connection.fetchval(
                "SELECT COUNT(*) FROM subscriptions WHERE plan = 'monthly' AND status = 'active'")
            yearly = await connection.fetchval(
                "SELECT COUNT(*) FROM subscriptions WHERE plan = 'yearly' AND status = 'active'")
            return users, monthly, yearly

db = Database(DATABASE_URL)

def get_main_keyboard(user_id):
    buttons = [
        [KeyboardButton(text="ПОМОЩЬ"), KeyboardButton(text="ПОДПИСКА")]
    ]
    if user_id == OWNER_CHAT_ID:
        buttons[0].append(KeyboardButton(text="АДМИН"))
    return ReplyKeyboardMarkup(keyboard=buttons, resize_keyboard=True)

@dp.message(CommandStart())
async def cmd_start(message: types.Message):
    await db.add_user(message.from_user.id)
    keyboard = get_main_keyboard(message.from_user.id)
    await message.answer(
        "Привет! Я твой BEST FRIEND 🤖\nГотов помочь с любыми вопросами!",
        reply_markup=keyboard
    )

@dp.message(F.text.lower() == "помощь")
async def help_command(message: types.Message):
    help_text = (
        "<b>🤖 BEST FRIEND: Правила и все возможности</b>\n\n"
        "👋 Привет! Я твой личный ИИ-бот для жизни, бизнеса, обучения, творчества, развлечения и не только. Вот что я умею:\n\n"
        "<b>⚙️ Мой движок (что под капотом):</b>\n\n"
        "    <b>GPT-4o</b> — топовая нейросеть от OpenAI, понимает и пишет на русском и английском, помнит твой диалог, гениальна почти во всём (в пределах разумного)!\n\n"
        "    <b>DALL-E 3</b> — нейросеть для генерации картинок. Можешь писать: <i>“нарисуй тигра в очках”</i> и получать уникальные изображения!\n\n"
        "    <b>GPT-4o Vision</b> — распознаёт и анализирует фото и картинки. Присылай фотки — расскажу, что на них.\n\n"
        "    <b>Whisper</b> — перевожу твои голосовые в текст. Хочешь поговорить — записывай аудио.\n\n"
        "<b>📄 Что можно отправлять:</b>\n\n"
        "    📝 <b>Текстовые запросы</b> (диалоги, статьи, переводы, стихи, идеи, коды и т.д.)\n\n"
        "    🎤 <b>Голосовые сообщения</b> — всё пойму, переведу в текст, отвечу.\n\n"
        "    🖼️ <b>Картинки и фото</b> — пришли, и я опишу, что на них, или сгенерирую новые по твоему описанию.\n\n"
        "    📚 <b>Файлы и документы</b> — поддерживаются PDF, Word, Excel, PowerPoint, txt, csv, zip, rar и многие другие! Я читаю, конвертирую и анализирую (до 4000 символов из файла).\n\n"
        "<b>🚦 Как пользоваться:</b>\n\n"
        "    <b>Пиши или говори</b> — просто формулируй любой вопрос, просьбу или задачу. Например, - \"составь курс обучения по...\"\n\n"
        "    <b>Для картинок</b> — начинай с «нарисуй», «создай», «сделай картинку…» (или пиши по-английски: generate, draw…).\n\n"
        "    <b>Для файлов</b> — просто отправляй документ, бот сам распознает тип и вытащит текст.\n\n"
        "<b>📎 Форматы ответов:</b>\n\n"
        "    Если твой ответ большой или содержит код — прилетит отдельным файлом.\n\n"
        "    Голосовые преобразую в текст и отвечаю тут же.\n\n"
        "    На картинках отправляю фото/иллюстрацию.\n\n"
        "    На файлах — распознаю и даю текстовое резюме.\n\n"
        "<b>🚫 <u>Правила использования:</u></b>\n\n"
        "    Не флуди — лимит запросов может быть ограничен по подписке.\n\n"
        "    Не отправляй запрещёнку: порнографию, угрозы, наркотики, попытки взлома, спам и т.д.\n\n"
        "    Уважай других — бот не поддерживает оскорбления и травлю.\n\n"
        "    Используй возможности для пользы, жизни, бизнеса, учёбы, личностного роста.\n\n"
        "    <b>Подписка</b> даёт расширенный доступ и отсутствие ограничений. Оформить/продлить можно через кнопку ПОДПИСКА.\n\n"
        "<b>🤪 Немного юмора:</b>\n"
        "Не стесняйся экспериментировать — я не психотерапевт, но могу посочувствовать, рассмешить, подбодрить и даже нарисовать слона в шляпе или шляпу слона.\n"
        "Если что-то не работает — просто напиши “почему ты тупишь?” или “админ”, и тебя поймут :)\n\n"
        "<b>💡 Запомни:</b>\n"
        "Я здесь, чтобы делать твою жизнь проще, интереснее и продуктивнее. Всё, что ты пишешь, обрабатывается ИИ. Не перегружай бот задачами на создание мирового господства — у меня лимит по электросети 😏\n\n"
        "<b>👑 Все вопросы, пожелания и предложения — всегда открыт к апгрейду! Я обучаюсь в процессе</b>\n\n"
        "Если что-то упустил — просто попробуй, и ты увидишь, на что я способен!\n"
        "Твой <b>BEST FRIEND</b> 🤖"
    )
    await message.answer(help_text, reply_markup=get_main_keyboard(message.from_user.id))

@dp.message(F.text.lower() == "подписка")
async def sub_command(message: types.Message):
    sub_url = "https://your-payment-link.com"
    user_id = message.from_user.id
    sub = await db.get_user_subscription(user_id)
    if sub is None:
        text = (
            "😔 <b>Подписка не активна</b>\n\n"
            "Оформи подписку, чтобы получить расширенный доступ:\n"
            f"<a href=\"{sub_url}\">Оплатить подписку</a>"
        )
    else:
        expires = sub["expires_at"].strftime("%d.%m.%Y %H:%M")
        text = (
            f"✅ <b>Ваша подписка активна до: {expires}</b>\n\n"
            "Продлить подписку:\n"
            f"<a href=\"{sub_url}\">Оплатить подписку</a>"
        )
    await message.answer(
        text,
        disable_web_page_preview=True,
        reply_markup=get_main_keyboard(user_id)
    )

@dp.message(F.text.lower() == "админ")
async def admin_stats(message: types.Message):
    if message.from_user.id != OWNER_CHAT_ID:
        return
    users, monthly, yearly = await db.get_stats()
    msg = (
        f"<b>Статистика 👑</b>\n"
        f"Пользователей: <b>{users}</b>\n"
        f"Месячных подписок: <b>{monthly}</b>\n"
        f"Годовых подписок: <b>{yearly}</b>"
    )
    await message.answer(msg, reply_markup=get_main_keyboard(message.from_user.id))

@app.on_event("startup")
async def on_startup():
    await db.connect()
    logging.info("Database connected")
    await bot.delete_my_commands()

@app.on_event("shutdown")
async def on_shutdown():
    await db.disconnect()
    logging.info("Database disconnected")

@app.post("/webhook")
async def telegram_webhook(request: Request):
    body = await request.json()
    update = types.Update(**body)
    await dp.feed_update(bot, update)
    return {"ok": True}

IMAGE_KEYWORDS = [
    r"^(нарисуй|создай|сделай|сгенерируй)\s*(картинку|изображение)?",
    r"^(generate|draw|create|make)\s*(image|picture)?",
]

def is_web_search_query(text):
    KEYWORDS = [
        "новости", "сегодня", "что нового", "тренды", "актуальное", "произошло",
        "курс", "сколько стоит", "погода", "интернет", "последние события",
        "что в мире", "текущий", "запроси", "что случилось", "price", "weather", "latest", "now"
    ]
    return any(k in text.lower() for k in KEYWORDS)

def should_send_as_file(text):
    if re.search(r"```.*?```", text, re.DOTALL):
        return True
    if re.match(r"^\s*(def |class |import |from |#|\/\/|<\w+)", text.strip()):
        return True
    lines = text.strip().split("\n")
    if len(lines) > 8 and any(sym in lines[0] for sym in ("def", "class", "import", "from", "#", "//", "<")):
        return True
    if re.search(r"(def |class |import |from |#|\/\/|<\w+)", text):
        return True
    return False

async def generate_filename(prompt, answer):
    system_prompt = (
        "Ты помощник для Telegram. На входе описание задачи и текст ответа с кодом. "
        "Верни ТОЛЬКО короткое английское название файла (без лишнего текста), не более 3 слов, через нижнее подчёркивание, всегда с расширением .txt, "
        "пример: snake_game.txt, telegram_bot.txt, sql_export_script.txt"
    )
    response = openai_client.chat.completions.create(
        model="gpt-4o",
        messages=[
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": f"Запрос: {prompt}\nОтвет:\n{answer}\n"},
        ],
        max_tokens=20,
        temperature=0.1
    )
    name = response.choices[0].message.content.strip()
    if not name.endswith(".txt"):
        name = "answer.txt"
    return name

def is_time_question(text):
    text = text.lower()
    return bool(re.search(r"\b(время|час|time)\b", text))

@dp.message(F.text)
async def universal_image_handler(message: types.Message):
    await handle_text_or_image(message, message.text)

async def handle_text_or_image(message, text):
    user_id = message.from_user.id
    t = text.strip().lower()
    if t in ["помощь", "подписка", "админ"]:
        return

    # --- Web Search via OpenAI Responses API ---
    if is_web_search_query(text):
        await message.answer("🔎 Делаю поиск в интернете через OpenAI Web Search, подожди секунду...", reply_markup=get_main_keyboard(user_id))
        try:
            response = openai_client.responses.create(
                model="gpt-4.1",
                tools=[{
                    "type": "web_search_preview",
                    # Можно добавить "search_context_size": "medium", "user_location": {...}
                }],
                input=text
            )
            answer = response.output_text
            await message.answer(answer, reply_markup=get_main_keyboard(user_id))
        except Exception as e:
            await message.answer(f"Ошибка при поиске в интернете через OpenAI: {e}", reply_markup=get_main_keyboard(user_id))
        return
    # --- /Web Search ---

    for pattern in IMAGE_KEYWORDS:
        m = re.match(pattern, t)
        if m:
            desc = re.sub(pattern, '', text, count=1).strip(":,. \n")
            if not desc:
                await message.answer("Опиши, что нужно нарисовать 👩‍🎨", reply_markup=get_main_keyboard(user_id))
                return
            try:
                response = openai_client.images.generate(
                    model="dall-e-3",
                    prompt=desc,
                    n=1,
                    size="1024x1024"
                )
                image_url = response.data[0].url
                await message.answer_photo(image_url, caption="Готово! Если хочешь ещё — просто напиши новый запрос.", reply_markup=get_main_keyboard(user_id))
            except Exception as e:
                await message.answer("Ошибка при генерации картинки 😔", reply_markup=get_main_keyboard(user_id))
            return
    if is_time_question(text):
        now = datetime.datetime.now().strftime("%H:%M:%S")
        await message.answer(f"Сейчас {now}", reply_markup=get_main_keyboard(user_id))
        return

    await db.add_message(user_id, "user", text)
    history = await db.get_history(user_id, limit=16)
    try:
        gpt_response = openai_client.chat.completions.create(
            model="gpt-4o",
            messages=history,
        )
        answer = gpt_response.choices[0].message.content

        SEARCH_TRIGGERS = [
            "не имею доступа к текущему времени",
            "не имею доступа к интернету",
            "я не могу узнать",
            "я не знаю",
            "я не обладаю актуальной информацией",
            "моя база устарела",
            "не могу ответить на этот вопрос",
            "у меня нет информации",
            "по состоянию на",
            "пожалуйста, проверьте актуальную информацию",
            "моя база данных не обновляется в реальном времени",
            "вы можете посмотреть актуальную стоимость",
            "я не могу предоставить текущую цену"
        ]

        if any(x in answer.lower() for x in SEARCH_TRIGGERS):
            answer = "Я не нашёл свежей информации по твоему запросу."

        await db.add_message(user_id, "assistant", answer)

        if should_send_as_file(answer):
            file_name = await generate_filename(text, answer)
            with open(file_name, "w", encoding="utf-8") as f:
                f.write(answer)
            with open(file_name, "rb") as f:
                await message.answer_document(types.BufferedInputFile(f.read(), file_name), caption="Готово! Вот твой файл 👇", reply_markup=get_main_keyboard(user_id))
            os.remove(file_name)
        else:
            await message.answer(answer, reply_markup=get_main_keyboard(user_id))
    except Exception:
        await message.answer("Ошибка при получении ответа от ИИ 🤖", reply_markup=get_main_keyboard(user_id))

# ------- Голосовые сообщения (Whisper + GPT-4o + генерация картинок) --------
@dp.message(F.voice)
async def handle_voice(message: types.Message):
    user_id = message.from_user.id
    file_id = message.voice.file_id
    file = await bot.get_file(file_id)
    ogg_path = f"voice_{user_id}.ogg"
    wav_path = f"voice_{user_id}.wav"
    await bot.download_file(file.file_path, ogg_path)
    try:
        audio = AudioSegment.from_file(ogg_path, format="ogg")
        audio.export(wav_path, format="wav")
    except Exception:
        await message.answer("Не получилось обработать голосовое 😢", reply_markup=get_main_keyboard(user_id))
        return
    try:
        with open(wav_path, "rb") as audio_file:
            transcript = openai_client.audio.transcriptions.create(
                model="whisper-1",
                file=audio_file,
                response_format="text",
                language="ru"
            )
        user_text = transcript.text if hasattr(transcript, "text") else str(transcript)
    except Exception:
        await message.answer("Ошибка при распознавании голосового 😔", reply_markup=get_main_keyboard(user_id))
        return
    finally:
        try:
            os.remove(ogg_path)
            os.remove(wav_path)
        except Exception:
            pass
    await handle_text_or_image(message, user_text)

# ------- Распознавание изображений (GPT-4o Vision) --------
@dp.message(F.photo)
async def handle_photo(message: types.Message):
    pass  # Можно вставить сюда распознавание фото через vision, если потребуется

# ------- Распознавание документов (GPT-4o + резюме) --------
@dp.message(F.document)
async def handle_document(message: types.Message):
    user_id = message.from_user.id
    doc = message.document
    filename = doc.file_name.lower()
    file = await bot.get_file(doc.file_id)
    f = BytesIO()
    await bot.download_file(file.file_path, destination=f)
    f.seek(0)

    text = ""
    error = None
    try:
        if filename.endswith('.pdf'):
            pdf = PdfReader(f)
            text = "".join([page.extract_text() or "" for page in pdf.pages])
        elif filename.endswith('.docx'):
            docx_file = DocxDocument(f)
            text = "\n".join([p.text for p in docx_file.paragraphs])
        elif filename.endswith('.txt'):
            text = f.read().decode('utf-8', errors='ignore')
        elif filename.endswith('.csv'):
            f.seek(0)
            lines = []
            reader = csv.reader(f.read().decode('utf-8', errors='ignore').splitlines())
            for row in reader:
                lines.append(','.join(row))
            text = "\n".join(lines)
        elif filename.endswith('.xlsx'):
            wb = openpyxl.load_workbook(f, read_only=True)
            ws = wb.active
            lines = []
            for row in ws.iter_rows(values_only=True):
                lines.append("\t".join([str(cell) for cell in row if cell is not None]))
            text = "\n".join(lines)
        elif filename.endswith('.xls'):
            book = xlrd.open_workbook(file_contents=f.read())
            sheet = book.sheet_by_index(0)
            lines = []
            for rx in range(sheet.nrows):
                lines.append("\t".join([str(cell) for cell in sheet.row_values(rx)]))
            text = "\n".join(lines)
        elif filename.endswith('.pptx'):
            ppt = pptx.Presentation(f)
            slides = []
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slides.append(shape.text)
            text = "\n".join(slides)
        elif filename.endswith('.zip'):
            z = zipfile.ZipFile(f)
            file_list = z.namelist()
            text = "ZIP-файл, содержит:\n" + "\n".join(file_list)
        elif filename.endswith('.rar'):
            r = rarfile.RarFile(f)
            file_list = r.namelist()
            text = "RAR-файл, содержит:\n" + "\n".join(file_list)
        else:
            try:
                text = textract.process(filename, input_stream=f).decode('utf-8', errors='ignore')
            except Exception:
                error = "Не удалось прочитать файл этим методом."
    except Exception as e:
        error = f"Ошибка при чтении файла: {e}"

    if text:
        # Обрезаем до 4000 символов (GPT-4o ограничение в prompt)
        chunk = text[:4000]
        prompt = (
            f"Сделай краткое, структурированное резюме по этому тексту (выдели основные моменты, сохрани факты, пиши лаконично):\n\n{chunk}"
        )
        try:
            gpt_response = openai_client.chat.completions.create(
                model="gpt-4o",
                messages=[
                    {"role": "system", "content": "Ты профессиональный ассистент, делаешь структурированные краткие резюме по тексту документа, не придумываешь фактов."},
                    {"role": "user", "content": prompt}
                ],
                max_tokens=600
            )
            summary = gpt_response.choices[0].message.content.strip()
            await message.answer(
                f"📄 <b>Файл:</b> <i>{doc.file_name}</i>\n\n<b>Резюме документа:</b>\n{summary}",
                reply_markup=get_main_keyboard(user_id),
                parse_mode="HTML"
            )
        except Exception as e:
            await message.answer(f"Ошибка при резюмировании через GPT: {e}", reply_markup=get_main_keyboard(user_id))
    else:
        await message.answer(f"❌ Не удалось прочитать документ. {error or ''}",
                             reply_markup=get_main_keyboard(user_id))

if __name__ == "__main__":
    import uvicorn
    uvicorn.run("main:app", host="0.0.0.0", port=10000)















































